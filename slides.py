# -*- coding: utf-8 -*-
import copy
import win32com.client
from pptx import Presentation


# def duplicate_slide(pres, template):
#     blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts) - 1]
#     copied_slide = pres.slides.add_slide(blank_slide_layout)

#     for shp in template.shapes:
#         el = shp.element
#         newel = copy.deepcopy(el)
#         copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

#     for value in iter(template.part.rels):
#         # Make sure we don't copy a notesSlide relation as that won't exist
#         if "notesSlide" not in value.reltype:
#             copied_slide.part.rels.add_relationship(
#                 value.reltype,
#                 value._target,
#                 value.rId
#             )
#     return copied_slide


def duplicate_slide_win32(source: str, dest: str, count: int):
    ppt_instance = win32com.client.Dispatch('PowerPoint.Application')
    #open the powerpoint presentation headless in background
    read_only = True
    has_title = False
    window    = False
    prs = ppt_instance.Presentations.open(source, read_only, has_title, window)
    nr_slide = 1

    for insert_index in range(2, count + 2):
        prs.Slides(nr_slide).Copy()
        prs.Slides.Paste(Index=insert_index)

    prs.SaveAs(dest)
    prs.Close()
    #kills ppt_instance
    ppt_instance.Quit()
    del ppt_instance


def search_and_replace_placeholder(slide, placeholder, value):
    for holder in slide.placeholders:
        for para in holder.text_frame.paragraphs:
            for runs in para.runs:
                if runs.text == placeholder:
                    runs.text = value

    return slide
